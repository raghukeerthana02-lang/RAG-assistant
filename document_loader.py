from pathlib import Path

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def extract_text(file_path, filename=None):

    extension = Path(file_path).suffix.lower()
    filename = filename or Path(file_path).name

    if extension == ".pdf":
        return extract_pdf(file_path, filename)

    elif extension == ".docx":
        return extract_docx(file_path, filename)

    elif extension == ".pptx":
        return extract_pptx(file_path, filename)

    else:
        raise ValueError(f"Unsupported file type: {extension}")


def extract_pdf(file_path, filename):

    reader = PdfReader(file_path)

    pages = []

    for page_number, page in enumerate(reader.pages, start=1):

        extracted = page.extract_text()

        if extracted:

            pages.append(
                {
                    "filename": filename,
                    "page": page_number,
                    "text": extracted
                }
            )

    return pages


def extract_docx(file_path, filename):

    document = Document(file_path)

    text = ""

    for paragraph in document.paragraphs:
        text += paragraph.text + "\n"

    return [
        {
            "filename": filename,
            "page": 1,
            "text": text
        }
    ]


def extract_pptx(file_path, filename):

    presentation = Presentation(file_path)

    slides = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        text = ""

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text += shape.text + "\n"

        slides.append(
            {
                "filename": filename,
                "page": slide_number,
                "text": text
            }
        )

    return slides